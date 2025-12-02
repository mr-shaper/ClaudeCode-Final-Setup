import base64
import io
import logging
from typing import Optional

# Import parsing libraries
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)

def parse_document_content(media_type: str, base64_data: str) -> str:
    """
    Parse document content from base64 data based on media type.
    Returns extracted text or a placeholder message if parsing fails.
    """
    try:
        # Decode base64 data
        file_data = base64.b64decode(base64_data)
        file_stream = io.BytesIO(file_data)
        
        extracted_text = ""

        # 1. PDF Parsing
        if media_type == "application/pdf":
            if not pypdf:
                return "[Error: pypdf library not installed]"
            reader = pypdf.PdfReader(file_stream)
            text_parts = []
            for page in reader.pages:
                text_parts.append(page.extract_text())
            extracted_text = "\n".join(text_parts)

        # 2. Word Parsing (.docx)
        elif media_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            if not DocxDocument:
                return "[Error: python-docx library not installed]"
            doc = DocxDocument(file_stream)
            extracted_text = "\n".join([para.text for para in doc.paragraphs])

        # 3. Excel Parsing (.xlsx)
        elif media_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            if not openpyxl:
                return "[Error: openpyxl library not installed]"
            wb = openpyxl.load_workbook(file_stream, data_only=True)
            text_parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text_parts.append(f"--- Sheet: {sheet} ---")
                for row in ws.iter_rows(values_only=True):
                    # Convert row to CSV-like string
                    row_text = ",".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip(","): # Skip empty rows
                        text_parts.append(row_text)
            extracted_text = "\n".join(text_parts)

        # 4. PPT Parsing (.pptx)
        elif media_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            if not Presentation:
                return "[Error: python-pptx library not installed]"
            prs = Presentation(file_stream)
            text_parts = []
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            extracted_text = "\n".join(text_parts)

        else:
            return f"[Document: {media_type} (Unsupported format for text extraction)]"

        # Return extracted text with a header
        if extracted_text.strip():
            return f"[Document Content ({media_type})]:\n{extracted_text}"
        else:
            return f"[Document: {media_type} (Empty content or extraction failed)]"

    except Exception as e:
        logger.error(f"Error parsing document {media_type}: {str(e)}")
        return f"[Document: {media_type} (Error parsing content: {str(e)})]"
